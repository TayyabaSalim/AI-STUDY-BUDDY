import os

from langchain_ollama import OllamaLLM, OllamaEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_chroma import Chroma
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain_community.document_loaders import PyPDFLoader
from langchain_core.documents import Document

from pptx import Presentation


# ============================================================
# CONFIGURATION
# ============================================================

DB_DIR = "./db"

LLM_MODEL = "llama3.2"
EMBEDDING_MODEL = "nomic-embed-text"


# ============================================================
# LOCAL OLLAMA LLM
# ============================================================

llm = OllamaLLM(
    model=LLM_MODEL,
    temperature=0.1,
    num_predict=300
)


# ============================================================
# LOCAL OLLAMA EMBEDDINGS
# ============================================================

embeddings = OllamaEmbeddings(
    model=EMBEDDING_MODEL
)


# ============================================================
# POWERPOINT LOADER
# ============================================================

def load_powerpoint(file_path):

    presentation = Presentation(file_path)

    documents = []

    for slide_number, slide in enumerate(
        presentation.slides,
        start=1
    ):

        slide_text = []

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                text = shape.text.strip()

                if text:
                    slide_text.append(text)

        if slide_text:

            documents.append(
                Document(
                    page_content="\n".join(slide_text),
                    metadata={
                        "source": file_path,
                        "slide": slide_number
                    }
                )
            )

    return documents


# ============================================================
# LOAD DOCUMENT
# ============================================================

def load_document(file_path):

    extension = os.path.splitext(
        file_path
    )[1].lower()

    # PDF
    if extension == ".pdf":

        loader = PyPDFLoader(file_path)

        return loader.load()

    # PowerPoint
    elif extension == ".pptx":

        return load_powerpoint(file_path)

    # Text
    elif extension == ".txt":

        with open(
            file_path,
            "r",
            encoding="utf-8"
        ) as file:

            text = file.read()

        return [
            Document(
                page_content=text,
                metadata={
                    "source": file_path
                }
            )
        ]

    else:

        raise ValueError(
            "Unsupported file type. "
            "Please upload PDF, PPTX, or TXT."
        )


# ============================================================
# PROCESS NOTES
# ============================================================

def process_notes(file_path):

    try:

        # --------------------------------------------
        # Extract text
        # --------------------------------------------

        docs = load_document(file_path)

        if not docs:

            return "No readable text was found."

        # --------------------------------------------
        # Split text
        # --------------------------------------------

        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=900,
            chunk_overlap=100
        )

        split_docs = text_splitter.split_documents(
            docs
        )

        if not split_docs:

            return "No text chunks were created."

        # --------------------------------------------
        # Store in Chroma
        # --------------------------------------------

        Chroma.from_documents(
            documents=split_docs,
            embedding=embeddings,
            persist_directory=DB_DIR,
            collection_name="study_notes"
        )

        extension = os.path.splitext(
            file_path
        )[1].upper()

        return (
            f"Successfully processed {extension} file. "
            f"Created {len(split_docs)} study chunks."
        )

    except Exception as error:

        print("\n===== PROCESSING ERROR =====")
        print(error)
        print("============================\n")

        raise


# ============================================================
# GET VECTOR STORE
# ============================================================

def get_vector_store():

    if not os.path.exists(DB_DIR):

        return None

    return Chroma(
        persist_directory=DB_DIR,
        embedding_function=embeddings,
        collection_name="study_notes"
    )


# ============================================================
# GENERATE SUMMARY
# ============================================================

def generate_summary():

    vector_store = get_vector_store()

    if vector_store is None:

        return "Please upload and process your notes first."

    try:

        # Retrieve only the most relevant chunks.
        # Fewer chunks = less work for the local LLM.

        docs = vector_store.similarity_search(
            "main topics key concepts definitions important facts",
            k=4
        )

        if not docs:

            return "No relevant study material was found."

        context = "\n\n".join(
            doc.page_content
            for doc in docs
        )

        # --------------------------------------------
        # Short, focused prompt
        # --------------------------------------------

        prompt = ChatPromptTemplate.from_template(
            """
You are a study assistant.

Summarize ONLY the study material below.

Give 5 to 8 concise bullet points.

Include:
- Important concepts
- Key definitions
- Important facts
- Rules or formulas if present

Do not add outside information.

Study material:
{context}

Summary:
"""
        )

        chain = (
            prompt
            | llm
            | StrOutputParser()
        )

        result = chain.invoke(
            {
                "context": context
            }
        )

        return result.strip()

    except Exception as error:

        print("\n===== SUMMARY ERROR =====")
        print(error)
        print("=========================\n")

        return f"Could not generate summary: {error}"


# ============================================================
# GENERATE FLASHCARDS
# ============================================================

def generate_flashcards():

    vector_store = get_vector_store()

    if vector_store is None:

        return "Please upload and process your notes first."

    try:

        docs = vector_store.similarity_search(
            "definitions concepts facts important information",
            k=4
        )

        if not docs:

            return "No relevant study material was found."

        context = "\n\n".join(
            doc.page_content
            for doc in docs
        )

        prompt = ChatPromptTemplate.from_template(
            """
You are a study assistant.

Create exactly 5 flashcards using ONLY
the study material below.

Use exactly this format:

Q: Question
A: Answer

Rules:
- Exactly 5 flashcards
- No numbering
- No introduction
- No conclusion
- No extra text
- Answers must come only from the material

Study material:
{context}

Flashcards:
"""
        )

        chain = (
            prompt
            | llm
            | StrOutputParser()
        )

        result = chain.invoke(
            {
                "context": context
            }
        )

        return result.strip()

    except Exception as error:

        print("\n===== FLASHCARD ERROR =====")
        print(error)
        print("===========================\n")

        return f"Could not generate flashcards: {error}"